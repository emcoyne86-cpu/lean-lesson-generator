"""
builders.py — Document builders for all 4 outputs.
All functions return BytesIO objects ready for st.download_button.
"""

import io
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.datavalidation import DataValidation
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn as pqn
from lxml import etree
import copy


# ══════════════════════════════════════════════════════════════════════════════
# DOCX BUILDER — Lean Lesson + Scaffolded Lesson
# ══════════════════════════════════════════════════════════════════════════════

# Colors
NAVY   = RGBColor(0x1F, 0x38, 0x64)
CORAL  = RGBColor(0xFF, 0x6B, 0x6B)
TEAL   = RGBColor(0x4E, 0xCE, 0xC8)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
AMBER  = RGBColor(0xF3, 0x9C, 0x12)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
LTGRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DKTEXT = RGBColor(0x1A, 0x1A, 0x2E)

# Population colors
ALL_BG   = RGBColor(0xD5, 0xF5, 0xE3)
ELL_BG   = RGBColor(0xD6, 0xEA, 0xF8)
SCD_BG   = RGBColor(0xE8, 0xDA, 0xEF)
ALL_TEXT = RGBColor(0x1E, 0x84, 0x49)
ELL_TEXT = RGBColor(0x19, 0x56, 0x99)
SCD_TEXT = RGBColor(0x6C, 0x34, 0x83)

TAG_COLORS = {
    "BACKGROUND":   CORAL,
    "OBJECTIVE":    TEAL,
    "KEY QUESTION": PURPLE,
    "VOCAB":        AMBER,
    "MODEL":        GREEN,
    "READ ALOUD":   RGBColor(0x85, 0x29, 0x9E),
    "PRACTICE":     RGBColor(0xE6, 0x7E, 0x22),
    "PROTOCOL":     RGBColor(0x16, 0xA0, 0x85),
    "CHECK":        RGBColor(0xC0, 0x39, 0x2B),
    "CLOSE":        RGBColor(0x21, 0x8C, 0x74),
    "HW":           RGBColor(0x7F, 0x8C, 0x8D),
}

COL_A = Inches(0.7)
COL_B = Inches(3.9)
COL_C = Inches(2.4)

POP_LABEL = {"ALL": "ALL STUDENTS", "ELL": "ELL", "SCD": "SCD"}
POP_BG    = {"ALL": ALL_BG,  "ELL": ELL_BG,  "SCD": SCD_BG}
POP_TEXT  = {"ALL": ALL_TEXT, "ELL": ELL_TEXT, "SCD": SCD_TEXT}


def _set_cell_bg(cell, rgb: RGBColor):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}")
    existing = tcPr.find(qn("w:shd"))
    if existing is not None:
        tcPr.remove(existing)
    tcPr.append(shd)


def _set_col_width(cell, width):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn("w:tcW")):
        tcPr.remove(old)
    tcW = OxmlElement("w:tcW")
    tcW.set(qn("w:w"), str(int(width.twips)))
    tcW.set(qn("w:type"), "dxa")
    tcPr.append(tcW)


def _cell_text(cell, text, bold=False, italic=False, color=None, size=10, align=WD_ALIGN_PARAGRAPH.LEFT):
    cell.text = ""
    p = cell.paragraphs[0]
    p.alignment = align
    run = p.add_run(text)
    run.bold = bold
    run.italic = italic
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color


def _make_header_table(doc, lesson_data):
    vocab = lesson_data.get("key_vocabulary", [])
    vocab_str = ",  ".join(vocab) if isinstance(vocab, list) else str(vocab)
    materials = lesson_data.get("materials", [])
    materials_str = "\n".join(f"• {m}" for m in materials) if isinstance(materials, list) else str(materials)

    meta = [
        ("Lesson Title",      lesson_data.get("title", "")),
        ("Curriculum",        f"{lesson_data.get('curriculum','')}  |  {lesson_data.get('unit','')}  |  {lesson_data.get('lesson','')}"),
        ("Grade / Subject",   f"Grade {lesson_data.get('grade','')}  |  {lesson_data.get('subject','')}"),
        ("Duration",          lesson_data.get("duration", "")),
        ("Standards",         "  ·  ".join(lesson_data.get("standards", []))),
        ("Learning Targets",  lesson_data.get("objective", "")),
        ("Key Vocabulary",    vocab_str),
        ("Materials",         materials_str),
    ]

    tbl = doc.add_table(rows=len(meta), cols=2)
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.LEFT

    for i, (label, val) in enumerate(meta):
        row = tbl.rows[i]
        _set_cell_bg(row.cells[0], NAVY)
        _cell_text(row.cells[0], label, bold=True, color=WHITE, size=9)
        _set_col_width(row.cells[0], Inches(1.0))
        _cell_text(row.cells[1], val, size=9)
        _set_col_width(row.cells[1], Inches(6.0))

    doc.add_paragraph().paragraph_format.space_after = Pt(4)


def _make_lesson_table(doc):
    """One single continuous table for the entire lesson."""
    tbl = doc.add_table(rows=0, cols=3)
    tbl.style = "Table Grid"
    tbl.alignment = WD_TABLE_ALIGNMENT.LEFT

    tblPr = tbl._tbl.tblPr
    tblLayout = OxmlElement("w:tblLayout")
    tblLayout.set(qn("w:type"), "fixed")
    tblPr.append(tblLayout)
    tblW = OxmlElement("w:tblW")
    tblW.set(qn("w:w"), str(int(Inches(7.0).twips)))
    tblW.set(qn("w:type"), "dxa")
    tblPr.append(tblW)

    return tbl


def _add_section_hdr_row(tbl, name, duration):
    """Navy full-width header row as first row of the section table."""
    row = tbl.add_row()
    # Merge all 3 cells into one
    row.cells[0].merge(row.cells[2])
    cell = row.cells[0]
    _set_cell_bg(cell, NAVY)
    _set_col_width(cell, Inches(7.0))
    cell.text = ""
    p = cell.paragraphs[0]
    r = p.add_run(f"{name.upper()}  —  {duration}")
    r.bold = True
    r.font.size = Pt(10)
    r.font.color.rgb = WHITE
    row.height_rule = 1  # exact
    row.height = Pt(16)


def _add_step_to_table(tbl, step, show_scaffold=False):
    """Add one step row (and any scaffold rows) to the section table."""
    row = tbl.add_row()
    tag     = step.get("tag", "")
    timing  = step.get("timing", "")
    content = step.get("content", "")
    tag_color = TAG_COLORS.get(tag.upper(), NAVY)

    # ── Col A: tag badge ─────────────────────────────────────────────
    a = row.cells[0]
    _set_cell_bg(a, tag_color)
    _set_col_width(a, COL_A)
    a.text = ""
    ap = a.paragraphs[0]
    ap.alignment = WD_ALIGN_PARAGRAPH.CENTER
    ar = ap.add_run(f"[{tag}]")
    ar.bold = True; ar.font.size = Pt(7.5); ar.font.color.rgb = WHITE

    # ── Col B+C: content (merged) or B|C split for KEY QUESTION ──────
    if tag.upper() == "KEY QUESTION":
        b = row.cells[1]
        _set_col_width(b, COL_B)
        b.text = ""
        bp = b.paragraphs[0]
        br = bp.add_run(content)
        br.font.size = Pt(9); br.font.italic = True

        c = row.cells[2]
        _set_col_width(c, COL_C)
        c.text = ""
        cp = c.paragraphs[0]
        cr1 = cp.add_run("Anticipated Response:")
        cr1.bold = True; cr1.font.size = Pt(8); cr1.font.color.rgb = TEAL
        anticipated = step.get("anticipated_response", "")
        if anticipated:
            cp.add_run("\n")
            cr2 = cp.add_run(anticipated)
            cr2.font.size = Pt(9); cr2.font.italic = True
    else:
        # Merge B + C → wide content cell
        b = row.cells[1]
        b.merge(row.cells[2])
        merged = row.cells[1]
        _set_col_width(merged, Inches(6.3))
        merged.text = ""
        mp = merged.paragraphs[0]
        if timing:
            mt = mp.add_run(f"(~{timing} min)  ")
            mt.font.size = Pt(8); mt.font.italic = True
            mt.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        mc = mp.add_run(content)
        mc.font.size = Pt(9)

    # ── Scaffold rows ─────────────────────────────────────────────────
    if show_scaffold:
        for scaf in step.get("scaffolds", []):
            sr = tbl.add_row()
            pop = scaf.get("population", "ALL")
            bg  = POP_BG.get(pop, ALL_BG)
            txt = POP_TEXT.get(pop, ALL_TEXT)

            sa = sr.cells[0]
            _set_cell_bg(sa, bg); _set_col_width(sa, COL_A)
            sa.text = ""
            sap = sa.paragraphs[0]; sap.alignment = WD_ALIGN_PARAGRAPH.CENTER
            sar = sap.add_run(POP_LABEL.get(pop, pop))
            sar.bold = True; sar.font.size = Pt(7.5); sar.font.color.rgb = txt

            sb = sr.cells[1]; _set_col_width(sb, COL_B)
            sb.text = ""
            sbp = sb.paragraphs[0]
            sbr1 = sbp.add_run(f"[{scaf.get('type','')}]  ")
            sbr1.bold = True; sbr1.font.size = Pt(9); sbr1.font.color.rgb = txt
            sbr2 = sbp.add_run(scaf.get("content", ""))
            sbr2.font.size = Pt(9)

            sc = sr.cells[2]; _set_col_width(sc, COL_C)
            _cell_text(sc, scaf.get("teacher_note", ""), size=8, italic=True)


def _build_doc(lesson_data: dict, label: str, show_scaffold: bool) -> bytes:
    doc = Document()
    for sec in doc.sections:
        sec.top_margin    = Inches(0.75)
        sec.bottom_margin = Inches(0.75)
        sec.left_margin   = Inches(0.75)
        sec.right_margin  = Inches(0.75)

    # Title
    tp = doc.add_paragraph()
    tr = tp.add_run(lesson_data.get("title", label))
    tr.bold = True; tr.font.size = Pt(16); tr.font.color.rgb = NAVY
    tp.paragraph_format.space_after = Pt(4)

    sp = doc.add_paragraph()
    sr = sp.add_run(label.upper())
    sr.bold = True; sr.font.size = Pt(9); sr.font.color.rgb = CORAL
    sp.paragraph_format.space_after = Pt(6)

    _make_header_table(doc, lesson_data)

    # One continuous table for the entire lesson
    tbl = _make_lesson_table(doc)
    for section_data in lesson_data.get("sections", []):
        _add_section_hdr_row(tbl, section_data["name"], section_data.get("duration", ""))
        for step in section_data.get("steps", []):
            _add_step_to_table(tbl, step, show_scaffold=show_scaffold)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def build_lean_docx(lesson_data: dict) -> bytes:
    return _build_doc(lesson_data, "Lean Lesson", show_scaffold=False)


def build_scaffolded_docx(lesson_data: dict) -> bytes:
    return _build_doc(lesson_data, "Scaffolded Lesson", show_scaffold=True)


# ══════════════════════════════════════════════════════════════════════════════
# XLSX BUILDER — Scaffold Selector
# ══════════════════════════════════════════════════════════════════════════════

def _xfill(hex_): return PatternFill("solid", fgColor=hex_)
def _xfont(hex_, bold=False, size=10, italic=False):
    return Font(color=hex_, bold=bold, size=size, italic=italic, name="Calibri")
def _xborder():
    s = Side(border_style="thin", color="CCCCDD")
    return Border(top=s, bottom=s, left=s, right=s)


def build_selector_xlsx(selector_rows: list[dict], lesson_data: dict) -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Scaffold Selector"

    NAVY_F  = _xfill("1F3864"); ALL_F = _xfill("1E8449"); ELL_F = _xfill("195699"); SCD_F = _xfill("6C3483")
    ALL_R   = _xfill("D5F5E3"); ELL_R = _xfill("D6EAF8"); SCD_R = _xfill("E8DAEF")
    ALT_ALL = _xfill("EAF9F0"); ALT_ELL = _xfill("EBF5FB"); ALT_SCD = _xfill("F4ECF7")

    headers = ["Step #", "Tag", "Step Description", "ALL", "ELL", "SCD", "Teacher Directive (optional)"]
    col_ws  = [7, 18, 68, 7, 7, 7, 44]

    for c, (h, w) in enumerate(zip(headers, col_ws), 1):
        cell = ws.cell(row=1, column=c, value=h)
        ws.column_dimensions[get_column_letter(c)].width = w
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        if c == 4:
            cell.fill = ALL_F; cell.font = _xfont("FFFFFF", bold=True)
        elif c == 5:
            cell.fill = ELL_F; cell.font = _xfont("FFFFFF", bold=True)
        elif c == 6:
            cell.fill = SCD_F; cell.font = _xfont("FFFFFF", bold=True)
        else:
            cell.fill = NAVY_F; cell.font = _xfont("FFFFFF", bold=True)
    ws.row_dimensions[1].height = 28

    # Lesson info in a merged cell above the table
    ws.insert_rows(1)
    title_cell = ws.cell(row=1, column=1, value=f"Scaffold Selector — {lesson_data.get('title','')}  |  Grade {lesson_data.get('grade','')}  |  {lesson_data.get('subject','')}")
    title_cell.fill = _xfill("1F3864"); title_cell.font = _xfont("FFFFFF", bold=True, size=12)
    title_cell.alignment = Alignment(horizontal="left", vertical="center")
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=8)
    ws.row_dimensions[1].height = 24

    # Data validation — Y/N dropdown for population columns only
    dv_pop = DataValidation(type="list", formula1='"Y,N"', allow_blank=True)
    ws.add_data_validation(dv_pop)

    for i, row in enumerate(selector_rows):
        r = i + 3  # offset for header rows
        alt = i % 2 == 1
        vals = [row["step_num"], row["tag"], row["content"],
                row.get("all",""), row.get("ell",""), row.get("scd",""), row.get("teacher_directive","")]
        for c, val in enumerate(vals, 1):
            cell = ws.cell(row=r, column=c, value=val)
            cell.border = _xborder()
            cell.alignment = Alignment(horizontal="left" if c > 2 else "center",
                                       vertical="top", wrap_text=True)
            cell.font = _xfont("1A1A2E", bold=(c == 2))
        ws.row_dimensions[r].height = 40

        # Color-code population cells (now cols 4, 5, 6)
        for c, (fill_yes, fill_alt) in enumerate([(ALL_R, ALT_ALL), (ELL_R, ALT_ELL), (SCD_R, ALT_SCD)], 4):
            ws.cell(row=r, column=c).fill = fill_alt if alt else fill_yes
            ws.cell(row=r, column=c).alignment = Alignment(horizontal="center", vertical="center")

        # Apply validations to cols 4, 5, 6
        for c in [4, 5, 6]:
            dv_pop.add(ws.cell(row=r, column=c))

    ws.freeze_panes = "A3"
    ws.sheet_view.showGridLines = False

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


def read_selector_xlsx(xlsx_bytes: bytes) -> list[dict]:
    """Parse a filled selector back into a list of row dicts."""
    wb = openpyxl.load_workbook(io.BytesIO(xlsx_bytes))
    ws = wb.active

    rows = []
    # Find the header row (contains "Step #")
    header_row = None
    for r in ws.iter_rows():
        for cell in r:
            if str(cell.value or "").strip() == "Step #":
                header_row = cell.row
                break
        if header_row:
            break

    if not header_row:
        return rows

    headers = [str(ws.cell(row=header_row, column=c).value or "").strip()
               for c in range(1, 9)]

    for r in range(header_row + 1, ws.max_row + 1):
        vals = [ws.cell(row=r, column=c).value for c in range(1, 8)]
        if not any(vals):
            continue
        all_val = str(vals[3] or "").strip()
        ell_val = str(vals[4] or "").strip()
        scd_val = str(vals[5] or "").strip()
        # Scaffold = Yes if any population is marked Y
        scaffold = "Yes" if any(v.upper() == "Y" for v in [all_val, ell_val, scd_val]) else "No"
        rows.append({
            "step_num":          vals[0],
            "tag":               str(vals[1] or "").strip(),
            "content":           str(vals[2] or "").strip(),
            "scaffold":          scaffold,
            "all":               all_val,
            "ell":               ell_val,
            "scd":               scd_val,
            "teacher_directive": str(vals[6] or "").strip(),
        })
    return rows


# ══════════════════════════════════════════════════════════════════════════════
# PPTX BUILDER — Student-Facing Slides
# ══════════════════════════════════════════════════════════════════════════════

SLIDE_W = PInches(10)
SLIDE_H = PInches(5.625)

BG_MAP = {
    "CORAL":  PRGBColor(0xFF, 0x6B, 0x6B),
    "TEAL":   PRGBColor(0x4E, 0xCE, 0xC8),
    "PURPLE": PRGBColor(0x9B, 0x59, 0xB6),
    "AMBER":  PRGBColor(0xF3, 0x9C, 0x12),
    "GREEN":  PRGBColor(0x27, 0xAE, 0x60),
    "PEACH":  PRGBColor(0xFF, 0xD9, 0xC8),
    "WHITE":  PRGBColor(0xFF, 0xFF, 0xFF),
    "NAVY":   PRGBColor(0x1F, 0x38, 0x64),
}
TEXT_WHITE = PRGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK  = PRGBColor(0x1A, 0x1A, 0x2E)

LIGHT_BG = {"PEACH", "WHITE", "AMBER"}


def _slide_bg(slide, color_name: str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    rgb = BG_MAP.get(color_name.upper(), BG_MAP["TEAL"])
    fill.fore_color.rgb = rgb
    return rgb


def _tb(slide, text, x, y, w, h, size=18, bold=False, color=None,
        align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txb


def _add_rounded_rect(slide, x, y, w, h, color: PRGBColor, radius=0.1):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE — use rounded via XML
        PInches(x), PInches(y), PInches(w), PInches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _text_color(bg_name: str) -> PRGBColor:
    return TEXT_DARK if bg_name.upper() in LIGHT_BG else TEXT_WHITE


def build_slides_pptx(slides_data: list[dict], lesson_data: dict) -> bytes:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    tc = _text_color  # alias

    for slide_info in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        bg_name = slide_info.get("bg_color", "TEAL")
        bg_rgb  = _slide_bg(slide, bg_name)
        txt_col = tc(bg_name)
        title_text = slide_info.get("title", "")
        bullets    = slide_info.get("bullets", [])
        featured   = slide_info.get("featured_text", "")
        stype      = slide_info.get("type", "content")
        notes_text = slide_info.get("speaker_notes", "")

        # Decorative accent oval (top-right)
        accent_rgb = PRGBColor(0xFF, 0xFF, 0xFF) if bg_name.upper() not in LIGHT_BG else PRGBColor(0x1F, 0x38, 0x64)
        oval_shape = slide.shapes.add_shape(9, PInches(8.8), PInches(0), PInches(2.5), PInches(2.5))
        oval_shape.fill.solid()
        oval_shape.fill.fore_color.rgb = accent_rgb
        oval_shape.fill.fore_color.transparency = 0.85
        oval_shape.line.fill.background()

        if stype == "title":
            # Big centered title
            _tb(slide, lesson_data.get("grade","") + " · " + lesson_data.get("subject","") + " · " + lesson_data.get("lesson",""),
                0.4, 0.3, 9.2, 0.5, size=13, bold=False, color=txt_col, align=PP_ALIGN.CENTER)
            _tb(slide, title_text, 0.4, 1.0, 9.2, 2.5, size=40, bold=True, color=txt_col, align=PP_ALIGN.CENTER)
            if lesson_data.get("curriculum"):
                _tb(slide, lesson_data["curriculum"], 0.4, 4.9, 9.2, 0.5, size=11, color=txt_col, align=PP_ALIGN.CENTER)

        elif stype == "target":
            _tb(slide, "Learning Target", 0.4, 0.2, 9.2, 0.55, size=14, bold=True, color=txt_col)
            # White card
            card = _add_rounded_rect(slide, 0.5, 0.9, 9.0, 3.8, PRGBColor(0xFF, 0xFF, 0xFF))
            card.fill.fore_color.transparency = 0.1
            _tb(slide, title_text, 0.7, 1.0, 8.6, 3.5, size=26, bold=True,
                color=PRGBColor(0x1A, 0x1A, 0x2E), align=PP_ALIGN.CENTER)
            if bullets:
                _tb(slide, "By the end of today, you will be able to:", 0.4, 4.85, 9.2, 0.5,
                    size=11, italic=True, color=txt_col)

        elif stype == "agenda":
            _tb(slide, title_text, 0.4, 0.2, 9.2, 0.55, size=22, bold=True, color=txt_col)
            y = 1.0
            for b in bullets[:6]:
                _tb(slide, "▸  " + b, 0.6, y, 8.8, 0.6, size=16, color=txt_col)
                y += 0.65

        elif stype in ("content", "materials"):
            _tb(slide, title_text, 0.4, 0.15, 9.2, 0.65, size=22, bold=True, color=txt_col)
            if featured:
                card = _add_rounded_rect(slide, 0.5, 0.95, 9.0, 2.4,
                                         PRGBColor(0xFF, 0xFF, 0xFF))
                card.fill.fore_color.transparency = 0.15
                _tb(slide, featured, 0.7, 1.0, 8.6, 2.3, size=20, bold=True,
                    color=TEXT_DARK, align=PP_ALIGN.CENTER)
                y = 3.5
            else:
                y = 0.95
            for b in bullets[:5]:
                _tb(slide, "▸  " + b, 0.6, y, 8.8, 0.65, size=15, color=txt_col)
                y += 0.7

        elif stype == "scaffold":
            # Peach background override for scaffold slides
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = PRGBColor(0xFF, 0xD9, 0xC8)
            txt_col = TEXT_DARK
            _tb(slide, title_text, 0.4, 0.15, 9.2, 0.65, size=20, bold=True, color=txt_col)
            if featured:
                card = _add_rounded_rect(slide, 0.5, 0.95, 9.0, 2.8, PRGBColor(0xFF, 0xFF, 0xFF))
                _tb(slide, featured, 0.7, 1.05, 8.6, 2.6, size=22, bold=True,
                    color=TEXT_DARK, align=PP_ALIGN.CENTER)
                y = 3.95
            else:
                y = 0.95
            for b in bullets[:4]:
                _tb(slide, "▸  " + b, 0.6, y, 8.8, 0.65, size=15, color=txt_col)
                y += 0.7

        elif stype == "closing":
            _tb(slide, title_text, 0.4, 0.2, 9.2, 0.65, size=22, bold=True, color=txt_col)
            if featured:
                card = _add_rounded_rect(slide, 0.5, 1.0, 9.0, 2.6, PRGBColor(0xFF, 0xFF, 0xFF))
                card.fill.fore_color.transparency = 0.1
                _tb(slide, featured, 0.7, 1.1, 8.6, 2.4, size=20, color=TEXT_DARK, align=PP_ALIGN.CENTER)
                y = 3.75
            else:
                y = 1.0
            for b in bullets[:4]:
                _tb(slide, "▸  " + b, 0.6, y, 8.8, 0.65, size=15, color=txt_col)
                y += 0.7

        # Speaker notes
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
